"""
RAG Chat — fastembed (local ONNX) + ChromaDB + Ollama Cloud LLM

Usage:
    python rag.py <folder_with_docs>

Example:
    python rag.py data

Type your question and press Enter. Type $quit to exit.

Environment Variables (.env):
    OLLAMA_API_KEY   - Your Ollama Cloud API key (required)
"""

import csv
import json
import os
import sys
import xml.etree.ElementTree as ET
from pathlib import Path

import chromadb
import fitz                      # pymupdf   — pip install pymupdf
import ollama
import openpyxl                  # pip install openpyxl
from chromadb.utils.embedding_functions import DefaultEmbeddingFunction
from docx import Document
from dotenv import load_dotenv
from pptx import Presentation    # pip install python-pptx

load_dotenv()

# ─── Configuration ────────────────────────────────────────────────────────────

OLLAMA_CLOUD_URL = "https://ollama.com"
LLM_MODEL = "llama3.1:8b"
EMBED_MODEL = "all-MiniLM-L6-v2"  # chromadb default ONNX model, downloads once
CHROMA_DIR = "./chroma_db"
CHUNK_SIZE = 500  # words
CHUNK_OVERLAP = 50  # words
TOP_K = 3

PROMPT = """\
Answer using ONLY the context below. Be concise.
If the answer is not in the context, say: "I couldn't find that in the documents."

Context:
{context}

Question: {question}
Answer:"""

# ─── Text splitting (no external deps) ───────────────────────────────────────


def split_text(
    text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP
) -> list[str]:
    words = text.split()
    chunks, i = [], 0
    step = size - overlap
    while i < len(words):
        chunks.append(" ".join(words[i : i + size]))
        i += step
    return chunks


# ─── Multi-format extraction ──────────────────────────────────────────────────


def extract_text(path: Path) -> str:
    ext = path.suffix.lower()

    if ext == ".txt":
        return path.read_text(encoding="utf-8", errors="ignore")

    if ext == ".md":
        return path.read_text(encoding="utf-8", errors="ignore")

    if ext == ".json":
        try:
            data = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
            return json.dumps(data, indent=2, ensure_ascii=False)
        except Exception:
            return path.read_text(encoding="utf-8", errors="ignore")

    if ext == ".csv":
        lines = []
        with open(path, newline="", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                lines.append(", ".join(row))
        return "\n".join(lines)

    if ext == ".xml":
        try:
            tree = ET.parse(path)
            return " ".join(el.text.strip() for el in tree.iter() if el.text and el.text.strip())
        except Exception:
            return path.read_text(encoding="utf-8", errors="ignore")

    if ext == ".docx":
        doc = Document(str(path))
        return "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())

    if ext == ".pdf":
        doc = fitz.open(str(path))
        pages = []
        for page in doc:
            pages.append(page.get_text())
        doc.close()
        return "\n".join(pages)

    if ext == ".xlsx":
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        rows = []
        for sheet in wb.worksheets:
            rows.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append(", ".join(cells))
        return "\n".join(rows)

    if ext == ".pptx":
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f"[Slide {i}]"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
            slides.append("\n".join(parts))
        return "\n\n".join(slides)

    print(f"  [SKIP] {path.name} — unsupported format ({ext})")
    return ""


def extract_structured_chunks(path: Path) -> list[dict]:
    ext = path.suffix.lower()

    # .docx: heading-aware chunking
    if ext == ".docx":
        doc = Document(str(path))
        chunks = []
        current_heading = "General"
        current_lines = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            if para.style.name.startswith("Heading"):
                if current_lines:
                    chunks.append({
                        "text": "\n".join(current_lines).lower(),
                        "heading": current_heading
                    })
                current_heading = text
                current_lines = []
            else:
                current_lines.append(text)
        if current_lines:
            chunks.append({
                "text": "\n".join(current_lines).lower(),
                "heading": current_heading
            })
        return chunks

    # All other formats: extract plain text, split into 300-word chunks
    raw = extract_text(path)
    if not raw.strip():
        return []
    words = raw.split()
    size, overlap = CHUNK_SIZE, CHUNK_OVERLAP
    step = size - overlap
    chunks = []
    i = 0
    while i < len(words):
        chunk_text = " ".join(words[i: i + size]).lower()
        chunks.append({"text": chunk_text, "heading": "General"})
        i += step
    return chunks


# ─── Helpers ──────────────────────────────────────────────────────────────────


def get_api_key() -> str:
    key = os.getenv("OLLAMA_API_KEY", "").strip()
    if not key:
        print("[ERROR] OLLAMA_API_KEY is not set in .env")
        sys.exit(1)
    return key


SUPPORTED_EXTENSIONS = {
    ".txt", ".md", ".json", ".csv", ".xml",
    ".docx", ".pdf", ".xlsx", ".pptx"
}


def load_and_index(folder: str) -> tuple[chromadb.Collection, list[Path]]:
    ef = DefaultEmbeddingFunction()
    ef(["warmup"])
    client = chromadb.PersistentClient(path=CHROMA_DIR)
    collection_name = Path(folder).resolve().name.replace(" ", "_")[:50] or "docs"

    all_files = sorted(
        p for p in Path(folder).iterdir()
        if p.is_file() and p.suffix.lower() in SUPPORTED_EXTENSIONS
    )

    # Return existing collection if already indexed
    existing = [c.name for c in client.list_collections()]
    if collection_name in existing:
        print(f"[INFO] Loading existing index '{collection_name}'...")
        return client.get_collection(name=collection_name, embedding_function=ef), all_files

    if not all_files:
        print(f"[ERROR] No supported files found in '{folder}'")
        sys.exit(1)

    print(f"[INFO] Indexing {len(all_files)} file(s)...")
    collection = client.create_collection(name=collection_name, embedding_function=ef)

    all_texts, all_ids, all_metas = [], [], []
    idx = 0
    for path in all_files:
        structured = extract_structured_chunks(path)
        if not structured:
            print(f"  [SKIP] {path.name} — no text extracted")
            continue
        for item in structured:
            all_texts.append(item["text"])
            all_ids.append(f"doc_{idx}")
            all_metas.append({
                "source": path.name,
                "heading": item["heading"]
            })
            idx += 1
        print(f"  [OK]   {path.name} — {len(structured)} chunks")

    if not all_texts:
        print("[ERROR] No text extracted.")
        sys.exit(1)

    # Add in batches of 100
    batch = 100
    for i in range(0, len(all_texts), batch):
        collection.add(
            documents=all_texts[i : i + batch],
            ids=all_ids[i : i + batch],
            metadatas=all_metas[i : i + batch],
        )
    print(f"[INFO] Done. {len(all_texts)} chunks stored.\n")
    return collection, all_files


def ask(question: str, collection, llm) -> None:
    results = collection.query(query_texts=[question], n_results=TOP_K)

    docs = results["documents"][0]
    metas = results["metadatas"][0]
    if not docs:
        print("\nBot: [No relevant context found]\n")
        return

    context = "\n\n---\n\n".join(
        f"[{m.get('source', '')}]\n{d}" for d, m in zip(docs, metas)
    )
    prompt = PROMPT.format(context=context, question=question)

    print("Bot: ", end="", flush=True)
    for chunk in llm.generate(model=LLM_MODEL, prompt=prompt, stream=True):
        print(chunk["response"], end="", flush=True)
    print("\n")


# ─── Main ─────────────────────────────────────────────────────────────────────


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    folder = sys.argv[1]
    if not os.path.isdir(folder):
        print(f"[ERROR] Not a directory: '{folder}'")
        sys.exit(1)

    api_key = get_api_key()
    collection, all_files = load_and_index(folder)

    llm = ollama.Client(
        host=OLLAMA_CLOUD_URL,
        headers={"Authorization": f"Bearer {api_key}"},
    )

    print("=" * 55)
    print(f"  {len(all_files)} document(s) ready:")
    for f in all_files:
        print(f"    - {f.name}")
    print("  Type $quit to exit.")
    print("=" * 55 + "\n")

    while True:
        try:
            question = input("You: ").strip()
        except (EOFError, KeyboardInterrupt):
            print("\nBye.")
            break
        if not question:
            continue
        if question == "$quit":
            print("Bye.")
            break
        ask(question, collection, llm)


if __name__ == "__main__":
    main()
